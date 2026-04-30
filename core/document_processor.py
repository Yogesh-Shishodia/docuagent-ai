"""Multi-format document parser.

Handles PDF, DOCX, XLSX/CSV, PPTX, EPUB, HTML, MD, TXT, JSON, and remote URLs.
Each parser returns a list of `Chunk` objects ready for embedding.
"""
from __future__ import annotations

import io
import json
import hashlib
import re
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Iterable

import pandas as pd

# Optional / heavy imports are done lazily inside the parsers
# to keep cold start fast and to avoid hard failures if a user
# doesn't need a given format.

from .config import CHUNK_SIZE, CHUNK_OVERLAP


# ---------------------------------------------------------------------------
# Data structures
# ---------------------------------------------------------------------------

@dataclass
class Chunk:
    """A single retrievable unit of a document."""
    text: str
    source: str          # original filename or URL
    doc_id: str          # stable hash for the source
    chunk_id: str        # unique within doc
    page: int | None = None
    section: str | None = None
    extra: dict = field(default_factory=dict)

    def to_metadata(self) -> dict:
        """Chroma-friendly metadata (must be flat scalar values)."""
        meta = {
            "source": self.source,
            "doc_id": self.doc_id,
            "chunk_id": self.chunk_id,
        }
        if self.page is not None:
            meta["page"] = int(self.page)
        if self.section:
            meta["section"] = str(self.section)[:200]
        return meta


def make_doc_id(name: str, content_sample: str) -> str:
    h = hashlib.sha1()
    h.update(name.encode("utf-8"))
    h.update(content_sample[:512].encode("utf-8", errors="ignore"))
    return h.hexdigest()[:12]


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------

def chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Sentence-aware sliding window chunker.

    Chunks at sentence boundaries when possible, falls back to hard splits.
    """
    text = re.sub(r"\s+", " ", text).strip()
    if len(text) <= size:
        return [text] if text else []

    # Split into sentences (cheap heuristic — good enough for chunking)
    sentences = re.split(r"(?<=[.!?])\s+(?=[A-Z0-9\"'(])", text)

    chunks: list[str] = []
    buf: list[str] = []
    buf_len = 0
    for sent in sentences:
        if buf_len + len(sent) + 1 <= size:
            buf.append(sent)
            buf_len += len(sent) + 1
        else:
            if buf:
                chunks.append(" ".join(buf).strip())
            # Build overlap from tail of previous buffer
            if overlap and chunks:
                tail = chunks[-1][-overlap:]
                buf = [tail, sent]
                buf_len = len(tail) + len(sent) + 1
            else:
                buf = [sent]
                buf_len = len(sent)
    if buf:
        chunks.append(" ".join(buf).strip())

    # Hard-split anything still too large
    final: list[str] = []
    for c in chunks:
        if len(c) <= size * 1.4:
            final.append(c)
        else:
            for i in range(0, len(c), size - overlap):
                final.append(c[i : i + size])
    return [c for c in final if c.strip()]


# ---------------------------------------------------------------------------
# Format-specific parsers
# ---------------------------------------------------------------------------

def parse_pdf(file_bytes: bytes, source: str) -> list[Chunk]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    pages_text: list[tuple[int, str]] = []
    for i, page in enumerate(reader.pages):
        try:
            pages_text.append((i + 1, page.extract_text() or ""))
        except Exception:
            pages_text.append((i + 1, ""))

    full_sample = " ".join(t for _, t in pages_text[:2])[:1000]
    doc_id = make_doc_id(source, full_sample)

    chunks: list[Chunk] = []
    for page_num, text in pages_text:
        for j, c in enumerate(chunk_text(text)):
            chunks.append(Chunk(
                text=c, source=source, doc_id=doc_id,
                chunk_id=f"{doc_id}-p{page_num}-{j}",
                page=page_num,
            ))
    return chunks


def parse_docx(file_bytes: bytes, source: str) -> list[Chunk]:
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    paragraphs: list[str] = []
    current_section = None
    section_chunks: list[tuple[str | None, str]] = []
    buf: list[str] = []

    for para in doc.paragraphs:
        if not para.text.strip():
            continue
        style = (para.style.name or "").lower()
        if "heading" in style:
            if buf:
                section_chunks.append((current_section, "\n".join(buf)))
                buf = []
            current_section = para.text.strip()
        else:
            buf.append(para.text.strip())
    if buf:
        section_chunks.append((current_section, "\n".join(buf)))

    # Tables — flatten as pipe-separated rows
    for table in doc.tables:
        rows = []
        for row in table.rows:
            rows.append(" | ".join(cell.text.strip() for cell in row.cells))
        if rows:
            section_chunks.append(("Table", "\n".join(rows)))

    full_text = "\n\n".join(t for _, t in section_chunks)
    doc_id = make_doc_id(source, full_text[:1000])

    chunks: list[Chunk] = []
    counter = 0
    for section, text in section_chunks:
        for c in chunk_text(text):
            chunks.append(Chunk(
                text=c, source=source, doc_id=doc_id,
                chunk_id=f"{doc_id}-d{counter}",
                section=section,
            ))
            counter += 1
    return chunks


def parse_xlsx(file_bytes: bytes, source: str) -> list[Chunk]:
    """Excel parsing — preserve sheet names and use a textual table format
    that the LLM can reason over. Also stash a small numeric profile for
    trend analysis later."""
    xls = pd.ExcelFile(io.BytesIO(file_bytes))
    chunks: list[Chunk] = []
    full_sample = ""

    for sheet in xls.sheet_names:
        df = xls.parse(sheet)
        if df.empty:
            continue
        df = df.dropna(how="all").dropna(how="all", axis=1)
        full_sample += df.head(3).to_string()[:500]

        # Profile — useful for trend analysis
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        profile = {
            "rows": int(len(df)),
            "columns": [str(c) for c in df.columns.tolist()],
            "numeric_columns": [str(c) for c in numeric_cols],
        }

        # Build a markdown-ish view of the sheet, chunked
        # Header line + a window of rows per chunk
        header = " | ".join(str(c) for c in df.columns)
        rows_per_chunk = 25
        for start in range(0, len(df), rows_per_chunk):
            slice_ = df.iloc[start : start + rows_per_chunk]
            body = "\n".join(
                " | ".join(str(v) for v in row) for row in slice_.itertuples(index=False)
            )
            text = f"Sheet: {sheet}\nRows {start+1}–{start+len(slice_)}\n{header}\n{body}"
            chunks.append(Chunk(
                text=text,
                source=source,
                doc_id="",  # filled below once we have enough sample
                chunk_id="",
                section=f"sheet:{sheet}",
                extra={"profile": profile},
            ))

    doc_id = make_doc_id(source, full_sample)
    for i, c in enumerate(chunks):
        c.doc_id = doc_id
        c.chunk_id = f"{doc_id}-x{i}"
    return chunks


def parse_csv(file_bytes: bytes, source: str) -> list[Chunk]:
    df = pd.read_csv(io.BytesIO(file_bytes))
    df = df.dropna(how="all").dropna(how="all", axis=1)
    full_sample = df.head(3).to_string()[:500]
    doc_id = make_doc_id(source, full_sample)

    profile = {
        "rows": int(len(df)),
        "columns": [str(c) for c in df.columns.tolist()],
        "numeric_columns": [str(c) for c in df.select_dtypes(include="number").columns.tolist()],
    }
    header = " | ".join(str(c) for c in df.columns)
    chunks: list[Chunk] = []
    rows_per_chunk = 25
    for start in range(0, len(df), rows_per_chunk):
        slice_ = df.iloc[start : start + rows_per_chunk]
        body = "\n".join(" | ".join(str(v) for v in r) for r in slice_.itertuples(index=False))
        text = f"Rows {start+1}–{start+len(slice_)}\n{header}\n{body}"
        chunks.append(Chunk(
            text=text, source=source, doc_id=doc_id,
            chunk_id=f"{doc_id}-c{start}",
            extra={"profile": profile},
        ))
    return chunks


def parse_pptx(file_bytes: bytes, source: str) -> list[Chunk]:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    slides: list[tuple[int, str]] = []
    for i, slide in enumerate(prs.slides, start=1):
        bits = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                bits.append(shape.text)
        slides.append((i, "\n".join(bits)))

    full_sample = " ".join(t for _, t in slides[:3])[:1000]
    doc_id = make_doc_id(source, full_sample)

    chunks: list[Chunk] = []
    for slide_num, text in slides:
        for j, c in enumerate(chunk_text(text)):
            chunks.append(Chunk(
                text=c, source=source, doc_id=doc_id,
                chunk_id=f"{doc_id}-s{slide_num}-{j}",
                page=slide_num, section=f"slide {slide_num}",
            ))
    return chunks


def parse_html(content: str, source: str) -> list[Chunk]:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(content, "lxml")
    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()
    text = soup.get_text("\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    return _generic_text_chunks(text, source)


def parse_url(url: str) -> list[Chunk]:
    """Fetch and extract main content from a URL using trafilatura."""
    import trafilatura
    downloaded = trafilatura.fetch_url(url)
    if not downloaded:
        raise ValueError(f"Could not fetch {url}")
    text = trafilatura.extract(downloaded, include_tables=True, include_comments=False)
    if not text:
        # Fallback to raw HTML parse
        return parse_html(downloaded, url)
    return _generic_text_chunks(text, url)


def parse_epub(file_bytes: bytes, source: str) -> list[Chunk]:
    from ebooklib import epub, ITEM_DOCUMENT
    from bs4 import BeautifulSoup

    book = epub.read_epub(io.BytesIO(file_bytes))
    sections: list[tuple[str, str]] = []
    for item in book.get_items_of_type(ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), "lxml")
        title = soup.find(["h1", "h2"])
        title_text = title.get_text(strip=True) if title else item.get_name()
        sections.append((title_text, soup.get_text("\n")))

    full_sample = " ".join(t for _, t in sections[:2])[:1000]
    doc_id = make_doc_id(source, full_sample)

    chunks: list[Chunk] = []
    counter = 0
    for sec_title, sec_text in sections:
        for c in chunk_text(sec_text):
            chunks.append(Chunk(
                text=c, source=source, doc_id=doc_id,
                chunk_id=f"{doc_id}-e{counter}",
                section=sec_title,
            ))
            counter += 1
    return chunks


def parse_text(content: str, source: str) -> list[Chunk]:
    return _generic_text_chunks(content, source)


def parse_json(content: str, source: str) -> list[Chunk]:
    try:
        obj = json.loads(content)
        pretty = json.dumps(obj, indent=2, ensure_ascii=False)
    except Exception:
        pretty = content
    return _generic_text_chunks(pretty, source)


def _generic_text_chunks(text: str, source: str) -> list[Chunk]:
    doc_id = make_doc_id(source, text[:1000])
    return [
        Chunk(text=c, source=source, doc_id=doc_id, chunk_id=f"{doc_id}-t{i}")
        for i, c in enumerate(chunk_text(text))
    ]


# ---------------------------------------------------------------------------
# Dispatch
# ---------------------------------------------------------------------------

def process_upload(name: str, data: bytes) -> list[Chunk]:
    """Route an uploaded file to the right parser based on extension."""
    ext = Path(name).suffix.lower().lstrip(".")
    try:
        if ext == "pdf":
            return parse_pdf(data, name)
        if ext == "docx":
            return parse_docx(data, name)
        if ext in ("xlsx", "xls"):
            return parse_xlsx(data, name)
        if ext == "csv":
            return parse_csv(data, name)
        if ext == "pptx":
            return parse_pptx(data, name)
        if ext == "epub":
            return parse_epub(data, name)
        if ext in ("html", "htm"):
            return parse_html(data.decode("utf-8", errors="replace"), name)
        if ext == "json":
            return parse_json(data.decode("utf-8", errors="replace"), name)
        # text-like fallthrough: txt, md, anything we can decode
        return parse_text(data.decode("utf-8", errors="replace"), name)
    except Exception as e:
        raise RuntimeError(f"Failed to parse {name}: {e}") from e


def document_summary_card(chunks: list[Chunk]) -> dict:
    """Lightweight statistics shown in the source list."""
    if not chunks:
        return {}
    pages = {c.page for c in chunks if c.page}
    sections = {c.section for c in chunks if c.section}
    total_chars = sum(len(c.text) for c in chunks)
    return {
        "doc_id": chunks[0].doc_id,
        "source": chunks[0].source,
        "chunks": len(chunks),
        "pages": len(pages) if pages else None,
        "sections": len(sections) if sections else None,
        "characters": total_chars,
        "approx_words": total_chars // 5,
    }